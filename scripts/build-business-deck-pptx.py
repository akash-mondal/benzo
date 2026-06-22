#!/usr/bin/env python3
"""Compile the Benzo-for-Business demo into a real PowerPoint (.pptx).

Per flow step: a SCREEN slide (screenshot + what-you-see + web2-clarity rating vs
Deel/Brex/Rippling/Ramp + advanced/on-chain) and a HOW-IT-WORKS slide (the
rendered ZK/contract diagram + contract-level + zero-knowledge-level deep text).
Plus title, mental model, the live M-of-N settle proof, scorecard, honesty, deploy.

Screens: docs/demo-flow/console-screens ; diagrams: docs/demo-flow/console-diagrams.
Writes docs/demo-flow/Benzo-Business-Demo.pptx.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SCREENS = os.path.join(ROOT, "docs/demo-flow/console-screens")
DIAGRAMS = os.path.join(ROOT, "docs/demo-flow/console-diagrams")
OUT = os.path.join(ROOT, "docs/demo-flow/Benzo-Business-Demo.pptx")

INK = RGBColor(0x19, 0x28, 0x37)
ACCENT = RGBColor(0x73, 0x42, 0xE2)
MUTED = RGBColor(0x6B, 0x6F, 0x74)
POS = RGBColor(0x1D, 0x7A, 0x52)
WARN = RGBColor(0x9A, 0x6B, 0x12)
CANVAS = RGBColor(0xF6, 0xF6, 0xF2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = CANVAS; bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    return tf


def para(tf, text, size, color, bold=False, first=False, space=6, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space)
    runs = text if isinstance(text, list) else [(text, color, bold)]
    for t, c, b in runs:
        r = p.add_run(); r.text = t; r.font.size = Pt(size); r.font.color.rgb = c
        r.font.bold = b; r.font.name = "Helvetica Neue"
    return p


def img_fit(s, path, x, y, max_w, max_h):
    if not os.path.exists(path):
        return
    from PIL import Image
    try:
        iw, ih = Image.open(path).size
    except Exception:
        iw, ih = 1440, 900
    r = min(max_w / iw, max_h / ih)
    s.shapes.add_picture(path, x + (max_w - int(iw * r)) // 2, y, int(iw * r), int(ih * r))


# ---- title ----
s = slide()
tf = tb(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(4.5))
para(tf, "Benzo for Business", 54, INK, bold=True, first=True, space=2)
para(tf, "Confidential payroll + AP, settled on-chain under M-of-N dual control", 24, ACCENT, bold=True, space=14)
para(tf, "Run payroll and pay vendors in USDC where every salary stays private on-chain, the treasury can only move "
        "with a member quorum (enforced in-circuit, not by a server), and an auditor can be given the total without the "
        "individual amounts. Zero-knowledge is invisible to the operator but load-bearing everywhere.", 15, MUTED)
para(tf, "Stellar / Soroban  ·  Groth16 over BN254  ·  Poseidon2  ·  in-circuit M-of-N (JSPLITORG, TEE-rotated VK)  ·  "
        "depth-32 Merkle pool  ·  Privacy-Pools ASP  ·  real Circle testnet USDC", 12, ACCENT, bold=True, space=2)

# ---- mental model ----
s = slide()
tf = tb(s, Inches(0.7), Inches(0.4), Inches(12), Inches(1.0))
para(tf, "The mental model — the chain is the backend, ZK is the security", 26, INK, bold=True, first=True)
para(tf, "The console is a desktop app over a client-side SDK. The treasury is a set of org notes owned by an M-of-N "
        "member set; every payout is a confidential transfer_org the Soroban verifier checks before a cent moves. "
        "No backend decides anything — the chain is the record.", 13.5, MUTED)
img_fit(s, os.path.join(DIAGRAMS, "diagram-0.png"), Inches(1.4), Inches(2.1), Inches(10.5), Inches(5.0))

STEPS = [
    dict(t="Sign up + KYB — verified on-chain", shot="01_onboarding_signup.png", dia="diagram-1.png", rate="9 / 10",
         see="SSO (no password), then a 6-step wizard: business details, KYB screening, compliance zone, team, treasury keys, review.",
         confused="A Deel/Rippling-grade onboarding; the only new idea is “treasury keys,” explained as “your secure books.”",
         adv="KYB is a real on-chain attestation in org_account (issuer-signed), not a backend flag; the console reads kyb_status from chain.",
         contract="org_account.register_org() records the org id, group key, threshold and members; attest_kyb() writes the screening decision signed by the issuer key. The console only READS kyb_status() — no backend adjudicates. SSO issues a session; no key or password is ever sent to a server.",
         zk="Sign-in derives the org viewing key client-side; the genuinely on-chain ZK step is registering the org’s authorized viewing key (MVK) in mvk_registry, which binds every future note to a member-readable scope without revealing balances."),
    dict(t="Fund the treasury — real USDC → a dual-controlled org note", shot="17_treasury_fund_prove.png", dia="diagram-2.png", rate="9 / 10",
         see="“Fund treasury” → amount → it lands as the treasury balance, badged “Provable.”",
         confused="Brex/Ramp “add funds” parity; the dual-control property is stated plainly, not jargon.",
         adv="Shields real USDC into an ORG note whose owner is a hash of (memberRoot, threshold, group key) — no single key can spend it. Live: BFF fund tx f4ba3185….",
         contract="pool.shield() pulls real testnet USDC and inserts a commitment whose recipient_pk = Poseidon2(memberRoot, threshold, akGroupPub; ORG). Because that owner is a hash preimage, the note can ONLY be moved by pool.transfer_org — never a single-key transfer.",
         zk="The shield proof binds the hidden amount to that org owner key. The treasury is then rediscovered from chain by decrypting the org notes with the org viewing key — no backend stores balances or member data."),
    dict(t="Roster + rate card (CSV) — amounts are computed", shot="11_contractors.png", dia="diagram-3.png", rate="9.5 / 10",
         see="Import a CSV (name, @handle, monthly USDC) or edit rates inline; “Run month” builds a payroll batch.",
         confused="This IS Deel/Gusto. Familiar to any HR/finance user; nothing crypto here.",
         adv="The engine computes each line from the rate card — it never sums numbers by hand. Each line becomes a confidential settle.",
         contract="No contract call at roster time — rate cards live in the org’s records. “Run month” computes line amounts and creates a payroll batch that enters the maker-checker flow.",
         zk="No proof yet; the computed amounts are private inputs to the confidential transfer_org settle for each contractor."),
    dict(t="Confidential payroll — M-of-N dual control, in circuit", shot="12_payroll.png", dia="diagram-4.png", rate="8.5 / 10",
         see="Approve & run a batch; each line settles a real shielded transfer. Individual salaries are never visible on-chain.",
         confused="Deel/Rippling “run payroll” parity. The privacy + dual-control are the upgrade, surfaced as plain status.",
         adv="Each payout is a pool.transfer_org settled under a >= threshold member quorum (JSPLITORG, TEE-rotated VK). Live: payouts 3a6de32c…, b9f87e4c….",
         contract="pool.transfer_org() verifies a Merkle root + the registered MVK root, spends the org nullifier, inserts the employee note + a fresh change org note, and calls verifier.verify_proof(JSPLITORG). The verifier REJECTS a single-key spend of org funds — release is gated by the proof inside the contract, not by the server.",
         zk="The joinsplit_org circuit (~147k constraints) proves: membership of the org note under a known root, value conservation, correct org nullifier, AND that >= threshold distinct members signed the spend message (EdDSA over Baby Jubjub). A sub-threshold or single-key spend cannot even produce a valid witness."),
    dict(t="Maker-checker — the approval gate IS the quorum", shot="13_approvals_maker_checker.png", dia="diagram-5.png", rate="9 / 10",
         see="A payment card shows “privacy by default” (amount + payee hidden), Approve/Deny, the approver trail, and “proposer can’t self-approve.”",
         confused="Brex/Ramp approval-inbox parity. Separation of duties is enforced, not just displayed.",
         adv="Each approval is one member signature over the exact spend; the release gate needs >= threshold before transfer_org will settle.",
         contract="The proposer (never an approver of their own payment) creates the order; approvers add signatures; at threshold the release gate fires pool.transfer_org. A consumer single-key transfer of org funds is rejected by the verifier — dual control is cryptographic.",
         zk="The maker-checker quorum is not a database flag: the M member signatures ARE the in-circuit witness the JSPLITORG proof checks. The on-chain verdict and the approval policy are the same fact."),
    dict(t="Pay a vendor / AP invoice — same engine, 2nd front door", shot="15_invoices_ap.png", dia="diagram-6.png", rate="9 / 10",
         see="An invoice inbox: pay one or pay all; over-threshold invoices route to approvals, under-threshold settle straight away.",
         confused="Bill.com / Ramp AP parity; status vocabulary (open / needs approval / paid) is familiar.",
         adv="Vendor pay settles through the same confidential transfer_org path; a double-entry ledger + payslip + CSV/GL export back it.",
         contract="pool.transfer_org() again — invoices are just a second front-door into the same settle engine. The approval policy decides whether a given invoice needs the maker-checker quorum.",
         zk="Identical confidentiality: amount + counterparty hidden on-chain, status plaintext. The audit ledger is a tamper-evident hash-chain over the settled records."),
    dict(t="Prove to an auditor — the total, never the salaries", shot="17c_treasury_disclose_total_result.png", dia="diagram-7.png", rate="opt-in",
         see="“Prove a balance” (holds >= X) and “Disclose exact total” give an auditor assurance; Auditor Grants issues a scoped viewing key.",
         confused="A power feature; the result is a plain “verified” badge. No competitor offers cryptographic disclosure at all.",
         adv="proof_of_sum / proof_of_balance are Groth16 proofs verified on-chain; a granted TVK is decrypt-only, scoped to a period, never a signer.",
         contract="verifier.verify_proof(SUM / BALANCE) confirms the statement on-chain (fail-closed). Auditor Grants hands out a time-scoped viewing key (TVK) derived one-way from the org MVK — read-only, revocable, never spend authority.",
         zk="The sum/balance circuits reveal only the aggregate (total = T, or holds >= X) with every individual note amount as a private witness. A scoped TVK lets an auditor passively decrypt only the in-scope notes — selective disclosure without ever exposing the member keys or the full book."),
    dict(t="Compliant edges — ASP admission + proof-of-innocence", shot="10_dashboard.png", dia="diagram-8.png", rate="—",
         see="Funds enter through an association allow-set and exit with a proof they’re not on a deny-set — private in the middle, accountable at the edges.",
         confused="Invisible to the operator; it’s the regulator answer behind the privacy.",
         adv="shield checks the asp_membership allow-root; withdraw requires a non-membership (proof-of-innocence) witness against the deny-set.",
         contract="asp_membership gates the shield edge (allow-root match); the non-membership SMT gates the withdraw edge. The middle (transfers) is fully private.",
         zk="Proof-of-innocence is a sparse-Merkle non-membership proof at the exit — the Privacy-Pools association-set model: fully private inside, provably clean on the way out."),
    dict(t="Console pays → contractor’s wallet → they spend it", shot="19_invites_roles.png", dia="diagram-9.png", rate="9 / 10",
         see="Invite a member / contractor / customer; they sign up in the consumer wallet and receive pay there, then send or cash out.",
         confused="Deel’s “invite a contractor” pattern; here the two apps actually interoperate over one shielded pool.",
         adv="The payout note is sealed to the contractor’s viewing key; their wallet rediscovers it from chain and can spend/cash-out. Verified at the SDK layer (employee rediscovers pay in a fresh client).",
         contract="The org’s transfer_org output note recipient is the contractor’s @handle key; nothing is custodial. The contractor’s wallet calls pool.transfer / pool.withdraw to move or cash out — the same pool, the other app.",
         zk="The two apps share the protocol, not a backend: the employee’s note is discoverable only with their viewing key, spendable only with their spend key. Console → wallet is a single confidential transfer with no link on-chain."),
    dict(t="Roles & permissions — separation of duties", shot="20_settings_roles_matrix.png", dia="diagram-10.png", rate="9 / 10",
         see="A roles matrix (owner / admin / treasurer / approver / auditor) shows exactly who can initiate, approve, release, and read.",
         confused="Brex/Rippling RBAC parity; the matrix makes the policy legible.",
         adv="Treasurer = release-gate signer; auditor = scoped viewing-key holder, never a signer. The roles map to the on-chain member set + threshold.",
         contract="The release gate maps to org_account’s threshold + member_root; only members of that root can contribute the signatures transfer_org requires.",
         zk="Org spends are unlinkable: the org nullifier nk_org = Poseidon2(akGroup, blinding) means two notes of the same set produce unrelated nullifiers — no org spend-graph leaks on-chain, while double-spends are still rejected by the nullifier set."),
]

for i, st in enumerate(STEPS):
    letter = chr(ord("a") + i)
    s = slide()
    tf = tb(s, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.7))
    para(tf, f"{letter}.  {st['t']}", 22, INK, bold=True, first=True)
    img_fit(s, os.path.join(SCREENS, st["shot"]), Inches(0.55), Inches(1.25), Inches(7.2), Inches(5.3))
    rtf = tb(s, Inches(8.0), Inches(1.3), Inches(4.9), Inches(5.5))
    para(rtf, [("What you see", ACCENT, True)], 13, ACCENT, first=True, space=2)
    para(rtf, st["see"], 12.5, INK, space=10)
    para(rtf, [(f"Web2-clarity vs Deel/Brex/Rippling/Ramp: {st['rate']}", POS, True)], 12, POS, space=2)
    para(rtf, st["confused"], 11.5, MUTED, space=10)
    para(rtf, [("Advanced · on-chain", ACCENT, True)], 11.5, ACCENT, space=2)
    para(rtf, st["adv"], 11.5, MUTED, space=2)
    # how-it-works
    s = slide()
    tf = tb(s, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.7))
    para(tf, f"{letter}.  {st['t']} — how it works", 21, INK, bold=True, first=True)
    img_fit(s, os.path.join(DIAGRAMS, st["dia"]), Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.7))
    rtf = tb(s, Inches(6.8), Inches(1.25), Inches(6.1), Inches(5.7))
    para(rtf, [("At the contract level", ACCENT, True)], 14, ACCENT, first=True, space=3)
    para(rtf, st["contract"], 12, INK, space=12)
    para(rtf, [("At the zero-knowledge level", ACCENT, True)], 14, ACCENT, space=3)
    para(rtf, st["zk"], 12, INK, space=2)

# ---- live M-of-N settle proof ----
s = slide()
tf = tb(s, Inches(0.7), Inches(0.45), Inches(12), Inches(6.6))
para(tf, "The moat, proven live on testnet — confidential payroll via M-of-N", 24, INK, bold=True, first=True, space=10)
para(tf, "Real Circle testnet USDC, the TEE-rotated JSPLITORG verification key:", 14, MUTED, space=12)
for line in [
    "Treasury funded as an M-of-N org note (recipient = hash of memberRoot, threshold, group key)",
    "Payout A — pool.transfer_org under a 2-of-3 quorum, org nullifier spent (tx 3a6de32c…)",
    "Payout B — spends the CHANGE org note under a different 2-of-3 quorum (tx b9f87e4c…): treasury stays dual-controlled across payouts",
    "A sub-threshold (1-of-3) payout is refused by dual control — funds cannot move on one key",
    "An employee rediscovered their pay in a fresh wallet client and withdrew it (real USDC out)",
]:
    para(tf, [("✓  ", POS, True), (line, INK, False)], 13.5, INK, space=8)
para(tf, "Verified by tests/e2e/org-payroll-settle.mjs + console-org-payroll.mjs; the member_root is published on-chain "
        "via org_account.set_member_root; the VK was rotated by a TEE-attested ceremony (quote UpToDate).", 12.5, MUTED, space=2)

# ---- scorecard ----
s = slide()
tf = tb(s, Inches(0.7), Inches(0.45), Inches(12), Inches(6.6))
para(tf, "Web2-clarity scorecard (vs Deel / Brex / Rippling / Ramp)", 24, INK, bold=True, first=True, space=12)
for sec, rate, note in [
    ("Sign up + KYB", "9", "Deel/Rippling-grade; KYB is real + on-chain"),
    ("Fund treasury", "9", "Brex/Ramp parity; dual-control stated plainly"),
    ("Roster + rate card (CSV)", "9.5", "This IS Deel/Gusto — fully familiar"),
    ("Run confidential payroll", "8.5", "Run-payroll parity; privacy + M-of-N are the upgrade"),
    ("Maker-checker approvals", "9", "Brex/Ramp inbox; separation of duties enforced"),
    ("AP invoices / vendor pay", "9", "Bill.com/Ramp AP parity"),
    ("Auditor disclosure", "opt-in", "No competitor offers cryptographic disclosure at all"),
    ("Roles & permissions", "9", "Brex/Rippling RBAC parity, mapped to the on-chain quorum"),
]:
    para(tf, [(f"{sec}", INK, True), (f"   {rate}/10   ", POS, True), (note, MUTED, False)], 13, INK, space=7)

# ---- honesty ----
s = slide()
tf = tb(s, Inches(0.7), Inches(0.45), Inches(12), Inches(6.6))
para(tf, "What is REAL vs SIMULATED (honesty)", 24, INK, bold=True, first=True, space=12)
para(tf, [("REAL, on-chain, verified:  ", POS, True)], 14, POS, space=3)
para(tf, "in-circuit M-of-N dual control (joinsplit_org / transfer_org / JSPLITORG, TEE-rotated VK); real Circle testnet "
        "USDC custody; Poseidon2 commitments + depth-32 tree; org + consumer nullifiers; ASP admission + proof-of-innocence; "
        "on-chain KYB attestation in org_account; the member_root; proof_of_sum / proof_of_balance verified on-chain.", 13, INK, space=12)
para(tf, [("SIMULATED / labeled:  ", WARN, True)], 14, WARN, space=3)
para(tf, "the fiat charge/payout leg of the ramp (needs a licensed anchor); the KYB screening provider stand-in (the "
        "decision is real on-chain, the provider integration is the seam); SSO is a dev session (OAuth is a config swap). "
        "The auditor ZK org-sum over org notes is the next circuit (today the figure is disclosed via the org view key).", 13, INK, space=2)

# ---- deploy ----
s = slide()
tf = tb(s, Inches(0.7), Inches(0.45), Inches(12), Inches(6.6))
para(tf, "Deploy readiness", 24, INK, bold=True, first=True, space=12)
para(tf, "Desktop-only console over a client-side SDK: the blockchain is the backend, ZK is the security layer. The "
        "treasury is rediscovered from chain (no backend stores member data); org keys derive deterministically from the "
        "owner seed so the self-hosted and deployed apps interoperate.", 14, INK, space=12)
for line in [
    "@benzo/core + console-api build clean; the console pay path routes through transfer_org (M-of-N), not single-key.",
    "Treasury reads cached (fast); demo-mode fabrications removed (seeded balances + stub balance proof killed).",
    "Deploy = `vercel --prod` (SPA rewrites + /api → BFF proxy) + host the BFF; env files are safe to share (no secrets).",
]:
    para(tf, [("→  ", ACCENT, True), (line, INK, False)], 13, INK, space=9)

prs.save(OUT)
print("wrote", OUT)
