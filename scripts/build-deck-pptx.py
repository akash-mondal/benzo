#!/usr/bin/env python3
"""Compile the Benzo consumer-app demo into a real PowerPoint (.pptx).

Per flow step: a SCREEN slide (screenshot + what-you-see + web2-clarity rating +
advanced/on-chain) and a HOW-IT-WORKS slide (the rendered ZK/contract diagram +
contract-level + zero-knowledge-level deep text — no short lingo). Plus title,
mental model, two-user-verified, proving-times, scorecard, honesty, deploy.

Reads screenshots from docs/demo-flow/screens and diagrams from
docs/demo-flow/diagrams; writes docs/demo-flow/Benzo-Consumer-Demo.pptx.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SCREENS = os.path.join(ROOT, "docs/demo-flow/screens")
DIAGRAMS = os.path.join(ROOT, "docs/demo-flow/diagrams")
OUT = os.path.join(ROOT, "docs/demo-flow/Benzo-Consumer-Demo.pptx")

INK = RGBColor(0x19, 0x28, 0x37)
ACCENT = RGBColor(0x73, 0x42, 0xE2)
MUTED = RGBColor(0x6B, 0x6F, 0x74)
POS = RGBColor(0x1D, 0x7A, 0x52)
CANVAS = RGBColor(0xF6, 0xF6, 0xF2)
CARD = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = CANVAS; bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    return tf


def para(tf, text, size, color, bold=False, first=False, space=6, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space)
    runs = text if isinstance(text, list) else [(text, color, bold)]
    for t, c, b in runs:
        r = p.add_run(); r.text = t; r.font.size = Pt(size); r.font.color.rgb = c
        r.font.bold = b; r.font.name = "Helvetica Neue"
    return p


def chip(s, x, y, text, color, fill):
    box = s.shapes.add_shape(5, x, y, Inches(2.2), Inches(0.36))
    box.fill.solid(); box.fill.fore_color.rgb = fill; box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = color; r.font.name = "Helvetica Neue"


def img_fit(s, path, x, y, max_w, max_h):
    if not os.path.exists(path):
        return
    from PIL import Image
    try:
        iw, ih = Image.open(path).size
    except Exception:
        iw, ih = 800, 1600
    r = min(max_w / iw, max_h / ih)
    w, h = int(iw * r), int(ih * r)
    s.shapes.add_picture(path, x + (max_w - w) // 2, y + (max_h - h) // 2, w, h)


# ---- title ----
s = slide()
bar = s.shapes.add_shape(1, 0, 0, SW, Inches(0.18)); bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False
tf = tb(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.5))
para(tf, "Benzo", 54, INK, bold=True, first=True, space=2)
para(tf, "Private USDC, made to feel like Cash App", 26, ACCENT, bold=True, space=14)
para(tf, "A private-by-default USDC wallet on Stellar. Zero-knowledge is invisible to the user but "
        "load-bearing everywhere: every step is a real Groth16 proof verified on-chain.", 15, MUTED)
para(tf, "Stellar / Soroban  ·  Groth16 over BN254 (CAP-0074)  ·  Poseidon2 (CAP-0075)  ·  depth-32 Merkle pool  ·  "
        "Privacy-Pools ASP  ·  real Circle testnet USDC", 12, ACCENT, bold=True, space=2)

# ---- mental model ----
s = slide()
tf = tb(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.9))
para(tf, "The mental model — account → shielded pool → account", 26, INK, bold=True, first=True)
para(tf, "USDC enters the pool as a hidden commitment (a “note”), moves privately inside as join-splits, "
        "and exits as a public withdrawal. The browser holds the keys and builds the proof; the Soroban pool "
        "verifies it on Stellar’s native elliptic-curve host functions before moving a cent.", 14, MUTED)
img_fit(s, os.path.join(DIAGRAMS, "diagram-0.png"), Inches(1.2), Inches(2.0), Inches(11), Inches(5.0))

# ---- per-step content ----
STEPS = [
    dict(t="Sign up — client-side identity", shot="02_onboarding_auth.png", dia="diagram-1.png", rate="9.5 / 10",
         see="Tap to create a passkey (Face ID / fingerprint). No seed phrase, no email, no password.",
         confused="No web2 user is confused — it’s smoother than any web3 wallet and on par with Cash App sign-up.",
         adv="No on-chain write except an optional @handle in the handle_registry contract (your public phone-book entry).",
         contract="Signing up touches at most one contract: handle_registry. Its register() maps your chosen name (e.g. “alicepay”) to your address; resolve() turns a name back into a payable recipient. No contract ever receives your keys — there is no custodial account.",
         zk="The passkey (a WebAuthn credential bound to this device) yields a device secret. A key-derivation function splits it into three roles: a spend key (authorizes spends, used inside the proof circuits), a master viewing key / MVK (decrypt-only — reads your notes, can never move money), and a note-discovery key (finds incoming payments). Decrypt and spend are cryptographically separated."),
    dict(t="Add money — fiat ramp → shield", shot="04_cash_add_money.png", dia="diagram-2.png", rate="9 / 10",
         see="“Add money” → amount → it lands in seconds. A live “on-chain reserve” badge shows the ramp is real.",
         confused="“Add Cash” parity with Cash App; the reserve badge adds trust. Only the fiat charge is simulated.",
         adv="The ramp contract dispenses real USDC from its on-chain reserve; the shield then hides the amount. Verified live: $5 added, reserve 5.00→0.00, ~42s.",
         contract="ramp.cash_in() dispenses real testnet USDC from the contract’s reserve to your funding address (the on-chain analog of a MoneyGram / SEP-24 anchor distribution account). pool.shield() then takes custody of that USDC and inserts a new commitment leaf.",
         zk="The shield circuit proves the new note’s commitment = Poseidon2(amount, recipientPk, blinding) is well-formed and that the deposited USDC equals the hidden amount — without revealing the amount. The commitment is appended to the depth-32 Merkle tree; only you can later open it."),
    dict(t="Import USDC — deposit from another wallet", shot="07_deposit_import_qr.png", dia="diagram-3.png", rate="8.5 / 10",
         see="A QR + your Stellar USDC address. Send from any wallet/exchange; when it lands, one tap shields it.",
         confused="“Receive crypto” parity, friendlier. The address is the only web3 artifact, and it’s clearly labeled.",
         adv="The screen polls the chain and shows “Ready to import $X · Detected on-chain.” Shielding it makes it private + spendable.",
         contract="No custody change on receipt — the USDC simply arrives at your public Stellar address. “Shield it” calls pool.shield() to move that public USDC into a hidden note, exactly like the ramp path but sourced from your own deposit.",
         zk="Identical to the add-money shield: a commitment is created and inserted; the amount becomes private the moment it is shielded. Before that, the public deposit is visible on-chain (honest — privacy begins at the shield edge)."),
    dict(t="Send privately — the join-split", shot="09_send.png", dia="diagram-4.png", rate="9 / 10",
         see="Pick a @handle or address, enter an amount, watch a 3-phase ceremony (building → proving → confirmed).",
         confused="Venmo-grade. The ceremony reassures; a first-time-recipient warning guards typos.",
         adv="On-chain this is a pool.transfer() — nullifiers + new commitments. You cannot read “Alice paid Bob $1” from it. Verified live: tx bb138474…, prove 1.27s.",
         contract="pool.transfer() consumes input notes by publishing their nullifiers (so they can never be respent), inserts new output commitments for recipient and change, and calls verifier_groth16.verify_proof() before any state change. The recipient is resolved off-chain via handle_registry.resolve().",
         zk="The joinsplit circuit proves: you own the input notes (Merkle membership against a known root), inputs = outputs (value conservation), each nullifier = f(spend key, leaf) is correctly derived, and the output commitments bind the hidden amounts to the recipient key — all in zero knowledge. Amounts and parties never appear on-chain."),
    dict(t="Request money", shot="12_request.png", dia="diagram-5.png", rate="8.5 / 10",
         see="Create a request link/QR for a chosen amount; share it. The payer just taps and sends.",
         confused="Familiar request-link pattern; decline + partial-pay beat Venmo.",
         adv="The request is a signed link (no on-chain write until the payer sends); the payment is the same private join-split.",
         contract="A request is an off-chain signed intent carrying your @handle + amount. No contract call until the payer acts — then it settles as a normal pool.transfer() to you.",
         zk="No proof at request time. The eventual payment carries the full join-split proof; the request just pre-fills the recipient + amount for the payer."),
    dict(t="Invite a new user → they claim & receive", shot="13_invite.png", dia="diagram-6.png", rate="8.5 / 10",
         see="“Invite” → escrow an amount → share a link. The invited person opens it, signs up, and the money is theirs.",
         confused="The $cashtag-link pattern; a self-refund if unclaimed is a plus. Verified live with a brand-new identity.",
         adv="Invite escrows into a claim account; claim sweeps it then shields into the new user’s own note. Verified: B 0→$2 (tx cc020c38…), then B sent $1 back to A (tx bb138474…).",
         contract="The invite funds a one-time claim account (or escrow contract) with a real on-chain hold; the link carries an app-scoped, domain-separated claim secret. claim() sweeps the escrowed note out and re-shields it under the recipient’s key; unclaimed past the time-lock, the sender self-refunds.",
         zk="The claim is a spend of the escrow note (proven with the claim secret) followed by a fresh shield into the recipient’s note — so the new user ends up with a note only they can spend, with no link between sender and recipient on-chain."),
    dict(t="Cash out — unshield + proof-of-innocence", shot="05_cash_cashout.png", dia="diagram-7.png", rate="9 / 10",
         see="“Cash out” → amount → “unshielded privately → returned to the reserve → to your bank.”",
         confused="Mirror of Add money; same trust cues. Verified live via the TEE prover: $5 cash-out settled (balance 11→6, reserve 0→5).",
         adv="The withdraw proves the note is NOT on the deny-list (proof-of-innocence) — private in the middle, accountable at the edges. Limits enforced ($5–$2,500).",
         contract="pool.withdraw() (the user-facing “unshield”) verifies the proof, publishes the spent note’s nullifier, and releases public USDC to your address; ramp.cash_out() then pulls it back into the reserve (the anchor absorbs it; only the fiat payout is simulated).",
         zk="The unshield circuit proves note ownership AND non-membership in the Association Set Provider (ASP) deny-list — a sparse-Merkle proof-of-innocence, mandatory at the exit edge. This is the regulator answer: fully private inside the pool, provably clean on the way out."),
    dict(t="Prove your balance — selective disclosure", shot="22_share_proof.png", dia="diagram-8.png", rate="n/a (opt-in)",
         see="Choose a threshold, get a proof you hold at least that much — never the exact amount.",
         confused="A power feature, hidden until needed; the result is a plain “Verified on-chain” badge.",
         adv="Generated on-device (witness never leaves the browser) and verified on-chain by the same verifier. Measured: cold 14.2s, warm 5.7s.",
         contract="No state change — verifier_groth16.verify_proof() is called as a read against the registered proof_of_balance verification key, returning true/false. The device can also re-verify the proof on-chain itself (trustless).",
         zk="The proof_of_balance circuit proves your spendable notes sum to ≥ the chosen threshold, with the exact balance as a private witness. The threshold is the only public input; your real balance stays hidden."),
    dict(t="Two ways to prove — device vs enclave (mobile = TEE)", shot="19_proving_mobile.png", dia="diagram-9.png", rate="—",
         see="The device decides: capable desktops prove on-device; phones + weak desktops delegate to the attested enclave (TEE).",
         confused="The user never picks — an auto chip explains it (“proving on-device” / “delegating to the secure enclave”).",
         adv="Soundness is identical either way — the on-chain verifier is the only trust anchor. The browser verifies the enclave’s TDX attestation + pins its measurement before sending the witness.",
         contract="Both paths submit to the same verifier_groth16 contract; the contract cannot tell (or care) where the proof was produced. The device-policy (10 unit tests) routes mobile + low-power machines to TEE.",
         zk="On-device: snarkjs WebAssembly prover, witness never leaves the browser (25 MB key, cached). TEE: the witness is sent over an attested RA-TLS channel sealed to the Intel Trusted Domain Extensions (TDX) enclave; a spoofed enclave fails the quote check and gets nothing. Verified live: TEE-proved unshield settled on-chain."),
    dict(t="The on-chain receipt — advanced disclosure", shot="23_txdetail_onchain.png", dia="diagram-10.png", rate="—",
         see="Every payment has a receipt: a plain status timeline + a hidden-but-not-too-hidden “View receipt” to the chain explorer.",
         confused="Normal-wallet receipt UX; the chain-scan link is there for judges, out of the way for everyone else.",
         adv="Shows Reference, Privacy: Private, Proof: Verified on-chain, and a Stellar-Expert explorer link for the settlement tx.",
         contract="The receipt reads the settled transaction + its nullifier/commitment events; “View receipt” deep-links to the explorer for the on-chain tx so anyone can independently confirm settlement.",
         zk="What’s public on the explorer is exactly the privacy boundary: the nullifier (spent) and the new commitment (created) — never the amount or the counterparties. The receipt makes that boundary legible."),
]

for i, st in enumerate(STEPS):
    letter = chr(ord("a") + i)
    # screen slide
    s = slide()
    tf = tb(s, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.7))
    para(tf, f"{letter}.  {st['t']}", 23, INK, bold=True, first=True)
    img_fit(s, os.path.join(SCREENS, st["shot"]), Inches(0.55), Inches(1.25), Inches(3.4), Inches(5.6))
    rtf = tb(s, Inches(4.3), Inches(1.3), Inches(8.4), Inches(5.4))
    para(rtf, [("What you see   ", ACCENT, True)], 13, ACCENT, first=True, space=2)
    para(rtf, st["see"], 14, INK, space=12)
    para(rtf, [(f"Web2-clarity: {st['rate']}", POS, True)], 13, POS, space=2)
    para(rtf, st["confused"], 12, MUTED, space=12)
    para(rtf, [("Advanced · on-chain   ", ACCENT, True)], 12, ACCENT, space=2)
    para(rtf, st["adv"], 12, MUTED, space=2)
    # how-it-works slide
    s = slide()
    tf = tb(s, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.7))
    para(tf, f"{letter}.  {st['t']} — how it works", 22, INK, bold=True, first=True)
    img_fit(s, os.path.join(DIAGRAMS, st["dia"]), Inches(0.5), Inches(1.2), Inches(5.7), Inches(5.6))
    rtf = tb(s, Inches(6.5), Inches(1.25), Inches(6.4), Inches(5.6))
    para(rtf, [("At the contract level", ACCENT, True)], 14, ACCENT, first=True, space=3)
    para(rtf, st["contract"], 12, INK, space=12)
    para(rtf, [("At the zero-knowledge level", ACCENT, True)], 14, ACCENT, space=3)
    para(rtf, st["zk"], 12, INK, space=2)

# ---- two-user verified ----
s = slide()
tf = tb(s, Inches(0.7), Inches(0.45), Inches(12), Inches(6.6))
para(tf, "Two-user flow — verified live on testnet", 24, INK, bold=True, first=True, space=10)
para(tf, "Two genuinely separate shielded identities (A on :8791, B on :8792), real Circle testnet USDC:", 14, MUTED, space=12)
for line in [
    "A creates a $2 invite — A 8.0 → 6.0 (escrowed), onChain",
    "B opens the link + claims — B 0 → $2.00, swept then shielded into B’s own note (tx cc020c38…)",
    "A registers @handle “alicepay” on-chain (tx abaf1ec2…)",
    "B sends $1 → alicepay (shielded join-split) — A 6.0 → 7.0, B 2.0 → 1.0 (tx bb138474…, prove 1.27s)",
]:
    para(tf, [("✓  ", POS, True), (line, INK, False)], 14, INK, space=8)
para(tf, "On-chain the B→A leg is a join-split, not “B paid A” — amounts and parties stay private.", 13, MUTED, space=2)

# ---- proving times ----
s = slide()
tf = tb(s, Inches(0.7), Inches(0.45), Inches(12), Inches(6.6))
para(tf, "Measured proving times (live testnet)", 24, INK, bold=True, first=True, space=10)
for label, val in [
    ("Browser on-device — COLD (note sync + 25 MB key load + WASM groth16)", "prove 12.5s + verify 1.7s ≈ 14.2s"),
    ("Browser on-device — WARM (key cached)", "prove 3.8s + verify 1.9s ≈ 5.7s"),
    ("Server prover (joinsplit, a real send)", "prove 1.27s (~3× faster than browser WASM)"),
    ("TEE enclave (Intel TDX) — cash-out unshield", "settled on-chain (attested path)"),
]:
    para(tf, [(label + "  —  ", INK, True), (val, ACCENT, True)], 14, INK, space=10)
para(tf, "Each verified on-chain (holds:true, onChain:true). The 25 MB key downloads once, then is cached "
        "(IndexedDB, encrypted at rest). The cold→warm drop is the one-time key load.", 13, MUTED, space=2)

# ---- scorecard ----
s = slide()
tf = tb(s, Inches(0.7), Inches(0.45), Inches(12), Inches(6.6))
para(tf, "Web2-clarity scorecard (vs Cash App / Venmo / Coinbase)", 24, INK, bold=True, first=True, space=12)
for sec, rate, note in [
    ("Sign up", "9.5", "No seed phrase or jargon — better than any web3 wallet"),
    ("Add money (ramp)", "9", "“Add Cash” parity; live reserve badge adds trust"),
    ("Import / deposit", "8.5", "“Receive crypto” parity, friendlier"),
    ("Send", "9", "Venmo-grade; ceremony reassures; new-recipient warning"),
    ("Request", "8.5", "Familiar; decline + partial-pay beat Venmo"),
    ("Invite / claim", "8.5", "$cashtag-link pattern; self-refund is a plus"),
    ("Cash out", "9", "Mirror of Add money"),
    ("Advanced details", "opt-in", "Hidden by default; one tap gives judges full chain provenance"),
]:
    para(tf, [(f"{sec}", INK, True), (f"   {rate}/10   ", POS, True), (note, MUTED, False)], 13, INK, space=7)

# ---- honesty ----
s = slide()
tf = tb(s, Inches(0.7), Inches(0.45), Inches(12), Inches(6.6))
para(tf, "What is REAL vs SIMULATED (honesty)", 24, INK, bold=True, first=True, space=12)
para(tf, [("REAL, on-chain, verified:  ", POS, True)], 14, POS, space=3)
para(tf, "every Groth16 proof verified on BN254; Poseidon2 commitments + tree; real Circle testnet USDC custody; "
        "nullifiers; ASP admission + proof-of-innocence; the ramp USDC legs; the on-chain @handle registry; escrow/claim.", 13, INK, space=12)
para(tf, [("SIMULATED:  ", RGBColor(0x9A, 0x6B, 0x12), True)], 14, RGBColor(0x9A, 0x6B, 0x12), space=3)
para(tf, "only the fiat charge/payout leg of the ramp (needs a licensed anchor) and the testnet KYC-tier stand-in "
        "(needs an identity-verification provider).", 13, INK, space=2)

# ---- deploy ----
s = slide()
tf = tb(s, Inches(0.7), Inches(0.45), Inches(12), Inches(6.6))
para(tf, "Deploy readiness", 24, INK, bold=True, first=True, space=12)
para(tf, "Client-side-first: the blockchain is the backend, zero-knowledge is the security layer. The browser holds "
        "the keys, builds the proof, and reads the chain; the only server-touch is a stateless gas relay + the "
        "handle directory (read) + the fiat ramp — none custodial.", 14, INK, space=12)
for line in [
    "Both apps prod-build clean; each ships a vercel.json (SPA rewrites + /api → BFF proxy); videos bundled.",
    "Unit suite 291 green (core 184 · kyc 24 · wallet 60 · console 23); cash-out + claim-shield bugs fixed.",
    "Deploy = `vercel --prod` per app + host the two BFFs; set BENZO_API_ORIGIN per the vercel.json rewrite.",
]:
    para(tf, [("→  ", ACCENT, True), (line, INK, False)], 13, INK, space=9)

prs.save(OUT)
print("wrote", OUT, "with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
